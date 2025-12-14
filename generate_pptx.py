from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Define slides content in structured format
slides = [
    {
        "title": "AI-Powered Migration Blueprint",
        "subtitle": "Accelerating Spring Boot to Node.js Transitions",
        "bullets": ["Innovation Showcase 2025", "December 14, 2025"],
        "style": "title"
    },
    {
        "title": "Use Case: What Problem Are We Solving?",
        "bullets": [
            "Teams struggle to migrate legacy Spring Boot apps to modern Node.js",
            "Manual migration is time-consuming and error-prone",
            "Maintaining functional parity is difficult",
            "Knowledge gaps between Java/Spring and JavaScript/Node ecosystems",
            "Delays, risk of bugs, inconsistent patterns, high experimentation costs"
        ],
        "style": "bullets"
    },
    {
        "title": "Approach: AI Tools & Methodology",
        "bullets": [
            "AI Tool: GitHub Copilot (Claude Sonnet 4.5)",
            "Intelligent scaffolding: Spring Boot layers, validation, tests",
            "Parallel migration: Node.js mirrors Spring controllers/services/repos",
            "Contract verification: OpenAPI, parity testing, identical seed data"
        ],
        "style": "bullets"
    },
    {
        "title": "Impact: Time Saved",
        "table": {
            "headers": ["Task", "Manual", "AI-Assisted", "Savings"],
            "rows": [
                ["Spring scaffold", "2-3 days", "15 minutes", "95%"],
                ["Node migration", "4-5 days", "20 minutes", "96%"],
                ["Test setup", "1-2 days", "10 minutes", "97%"],
                ["Documentation", "1 day", "5 minutes", "98%"],
                ["TOTAL", "8-11 days", "~50 minutes", "96%"]
            ]
        },
        "style": "table"
    },
    {
        "title": "Impact: Quality Improvements",
        "bullets": [
            "Zero syntax errors in generated code",
            "Consistent patterns across stacks (controller/service/repo)",
            "Production-ready validation, error handling, logging",
            "Complete test coverage structure (MockMvc, Jest/Supertest)",
            "Dockerfile + OpenAPI included"
        ],
        "style": "bullets"
    },
    {
        "title": "Creative Innovations",
        "bullets": [
            "Migration blueprint mono-repo template",
            "Living documentation via identical API contracts",
            "Zero-config paired demos",
            "Layered architecture mapping (controller/service/repo)"
        ],
        "style": "bullets"
    },
    {
        "title": "What We Built (Structure)",
        "bullets": [
            "packages/spring-app: controllers, services, repositories, DTOs, exceptions, Dockerfile, OpenAPI",
            "packages/node-migration: routes, services, repositories, validators, middleware, Dockerfile"
        ],
        "style": "bullets"
    },
    {
        "title": "Live Demo: API Parity",
        "bullets": [
            "Spring Boot (8080): GET /api/products → Widget, Gadget",
            "Node.js (4000): GET /api/products → Widget, Gadget",
            "Identical responses validate migration"
        ],
        "style": "bullets"
    },
    {
        "title": "Judging Criteria Alignment",
        "bullets": [
            "Practicality: solves a real migration problem",
            "Impact: ~96% time reduction; lower risk",
            "Creativity: mono-repo blueprint + contract-first",
            "Replicability: clean structure, minimal deps, forkable",
            "Presentation: clear metrics, demos, mapping"
        ],
        "style": "bullets"
    },
    {
        "title": "How Others Can Adopt",
        "bullets": [
            "Clone: https://github.com/kanchanakm8/spring-app",
            "Customize domain model (replace Product)",
            "Run side-by-side: Spring jar + Node start",
            "Compare outputs via curl to validate"
        ],
        "style": "bullets"
    },
    {
        "title": "ROI Calculation",
        "bullets": [
            "Traditional: $6.4K-$8.8K (8-11 days)",
            "AI-Assisted: ~$800 (1 day)",
            "Savings per project: $5.6K-$8K; 10 teams → $56K-$80K"
        ],
        "style": "bullets"
    },
    {
        "title": "Key Takeaways",
        "bullets": [
            "AI accelerates migration discovery by 20-50x",
            "Pattern-based translation bridges frameworks",
            "Mono-repo + contracts = reusable templates",
            "Small innovations compound value"
        ],
        "style": "bullets"
    },
    {
        "title": "Conclusion",
        "bullets": [
            "Weeks-long migration → 50-minute AI-assisted exercise",
            "Reusable blueprint; AI-powered translation; contract-verified parity",
            "Production-ready structure; 96% time reduction",
            "Call to action: adopt this blueprint"
        ],
        "style": "bullets"
    },
]

def add_title_slide(prs, slide):
    layout = prs.slide_layouts[0]  # Title
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = slide["title"]
    s.placeholders[1].text = "\n".join(slide.get("bullets", []))

def add_bullet_slide(prs, slide):
    layout = prs.slide_layouts[1]  # Title and Content
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = slide["title"]
    body = s.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, bullet in enumerate(slide.get("bullets", [])):
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0


def add_table_slide(prs, slide):
    layout = prs.slide_layouts[5]  # Title Only
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = slide["title"]
    rows = len(slide["table"]["rows"]) + 1
    cols = len(slide["table"]["headers"])
    table = s.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(8), Inches(4)).table
    # headers
    for c, h in enumerate(slide["table"]["headers"]):
        cell = table.cell(0, c)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
    # rows
    for r, row in enumerate(slide["table"]["rows"], start=1):
        for c, val in enumerate(row):
            table.cell(r, c).text = str(val)


def build_presentation(slides):
    prs = Presentation()
    for slide in slides:
        style = slide.get("style", "bullets")
        if style == "title":
            add_title_slide(prs, slide)
        elif style == "table":
            add_table_slide(prs, slide)
        else:
            add_bullet_slide(prs, slide)
    return prs


def main():
    prs = build_presentation(slides)
    output = "AI_Migration_Blueprint.pptx"
    prs.save(output)
    print(f"Saved {output}")

if __name__ == "__main__":
    main()
